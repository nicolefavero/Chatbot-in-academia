################################################################################
#  CHATBOT WITH OPENAI GPT-4o Mini PRODUCING SLIDE DECKS FROM ACADEMIC PAPERS
################################################################################
import os
import re
import json
import ast
import gradio as gr
from fuzzywuzzy import fuzz
from openai import OpenAI
from DOC_REGISTRY import DOC_REGISTRY
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT

client = OpenAI(api_key="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")

def preprocess_text(text: str) -> str:
    '''Function to clean and preprocess the text.
    Removes references, footnotes, and unnecessary characters.
    Args:
        text (str): The raw text to preprocess.
    Returns:
        str: The cleaned text.
    '''
    text = text.strip()
    text = re.sub(r"\[\d+\]|\(\w+ et al\., \d+\)", "", text)
    text = re.sub(r"\(see Fig\.\s?\d+\)", "", text)
    text = re.sub(r"[*_#]", "", text)
    text = re.sub(r"-\n", "", text)
    return " ".join(text.split())

def read_txt_full_text(file_path: str) -> str:
    '''Function to read the content of a text file.
    Args:
        file_path (str): The path to the text file.
    Returns:
        str: The content of the text file.'''
    with open(file_path, "r", encoding="utf-8") as f:
        raw_text = f.read()
    return preprocess_text(raw_text)

def detect_target_doc(query: str, registry: dict, threshold=80):
    '''Function to detect the target document based on the user's query.
    Args:
        query (str): The user's query.
        registry (dict): The document registry containing document names and aliases.
        threshold (int): The minimum score to consider a match. 
    Returns:
        str: The name of the matched document or None if no match is found.
    '''
    query_lower = query.lower()
    best_doc = None
    best_score = 0
    for doc_name, data in registry.items():
        for alias in data["aliases"]:
            score = fuzz.partial_ratio(query_lower, alias.lower())
            if score > best_score:
                best_score = score
                best_doc = doc_name
    return best_doc if best_score >= threshold else None

def gpt_response(system_content, user_prompt, max_tokens=1000):
    '''Function to get a response from the GPT model.
    Args:
        system_content (str): The system message for the model.
        user_prompt (str): The user's prompt.
        max_tokens (int): The maximum number of tokens to generate.
    Returns:
        str: The response from the model.
    '''
    response = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[
            {"role": "system", "content": system_content},
            {"role": "user", "content": user_prompt}
        ],
        temperature=0.7,
        max_tokens=max_tokens
    )
    return response.choices[0].message.content

def parse_python_list(raw_output: str) -> list:
    '''Function to parse a Python list from a string.
    Args:
        raw_output (str): The raw output string.
    Returns:
        list: The parsed Python list.
    '''
    try:
        bracketed = raw_output[raw_output.find("["):raw_output.rfind("]")+1]
        data = ast.literal_eval(bracketed)
        if isinstance(data, list):
            return data
    except Exception:
        pass
    return []

def generate_slide_deck(full_text: str):
    '''Function to generate a slide deck from the full text.
    Args:
        full_text (str): The full text of the paper.
    Returns:
        list: A list of slides with titles and bullet points.
    '''
    titles_prompt = f"""
You are an academic presentation assistant. Generate EXACTLY 5 concise slide titles as a Python list based on the following academic content:
Return only the list, no explanation.

Content:
{full_text[:3000]}
"""
    raw_titles = gpt_response("You generate slide titles only.", titles_prompt, 300)
    titles = parse_python_list(raw_titles)
    if len(titles) < 5:
        titles = [f"Slide {i+1}" for i in range(5)]

    slides = []
    for title in titles:
        bullets_prompt = f"""
You are an academic assistant. Generate EXACTLY 3 bullet points as a Python list for the slide titled '{title}' using the content below:
Return only the list.

Content:
{full_text[:3000]}
"""
        raw_bullets = gpt_response("You write 3 bullet points only.", bullets_prompt, 300)
        bullets = parse_python_list(raw_bullets)
        if len(bullets) < 3:
            bullets = ["Point 1", "Point 2", "Point 3"]
        slides.append({"title": title, "bullet_points": bullets})
    return slides

def build_pptx(slides, output_path="presentation.pptx"):
    '''Function to build a PowerPoint presentation from the slides.
    Args:
        slides (list): A list of slides with titles and bullet points.
        output_path (str): The path to save the PowerPoint file.
    Returns:
        str: The path to the saved PowerPoint file.
    '''
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]

    # Title Slide
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "Academic Slide Deck"
    slide.placeholders[1].text = "Generated by CBS Bot"

    # Content Slides
    for item in slides:
        slide = prs.slides.add_slide(bullet_slide_layout)
        slide.shapes.title.text = item["title"]
        content = slide.shapes.placeholders[1]
        tf = content.text_frame
        tf.clear()
        tf.word_wrap = True
        for point in item["bullet_points"]:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0
            p.font.size = Pt(16)  # Slightly smaller
            p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

    prs.save(output_path)
    return output_path

def handle_user_query(message, chat_history):
    '''Function to handle user queries and generate slide decks.
    Args:
        message (str): The user's message.
        chat_history (list): The chat history.
    Returns:
        str: The response to the user.
    '''
    matched_doc = detect_target_doc(message, DOC_REGISTRY)
    if not matched_doc:
        try:
            fallback_response = gpt_response(
                "You are a friendly assistant who only helps users generate slide decks for academic papers they name.",
                f"I couldn't find the paper '{message}'. Can you try a different title, or rephrase it?",
                200
            )
            return fallback_response
        except Exception as e:
            return f"⚠️ Error while handling unknown input: {e}"

    try:
        file_path = f"./papers-cleaned/{matched_doc}.txt"
        full_text = read_txt_full_text(file_path)

        slides = generate_slide_deck(full_text)
        pptx_path = build_pptx(slides)

        return gr.File(pptx_path, label="📥 Download your slides")

    except Exception as e:
        return f"⚠️ Error: {e}"

chatbot = gr.ChatInterface(
    fn=handle_user_query,
    title="📚 CBS Slide Generator",
    theme="default",
    chatbot=gr.Chatbot(show_copy_button=True, type="messages")
)

chatbot.launch()
