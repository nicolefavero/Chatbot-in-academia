################################################################################
#  ACADEMIC SUMMARY GENERATOR + SLIDE CREATOR (LLAMA 3 - 70B, TWO LLM CALLS)   #
################################################################################

import os
import re
import torch
from transformers import AutoModelForCausalLM, AutoTokenizer

# Import summary functions from LLama_Summary.py
from LLama_Summary import (
    load_llama_instructor,
    summarize_text_doc as llama_summarize_text_doc
)
from DOC_REGISTRY import DOC_REGISTRY

################################################################################
# 1. Summary Processing and Slide Content Generation
################################################################################

def extract_sections_from_summary(summary_text):
    """
    Extract structured sections from academic summary text using Markdown-style, bolded, or all-caps headers.
    """
    # Remove trailing helper signatures
    summary_text = re.split(r'(?:Best regards|Your Academic Writer|Note:)', summary_text)[0]

    # Match Markdown headings, bolded text, or ALL CAPS lines
    section_regex = re.compile(
        r'(?:^|\n)\s*(?:\*\*|#+)?\s*([A-Z][A-Za-z\s&\-]{3,60})\s*(?:\*\*|#+)?\s*\n(.*?)(?=(?:\n\s*(?:\*\*|#+)?\s*[A-Z][A-Za-z\s&\-]{3,60}\s*(?:\*\*|#+)?\s*\n)|\Z)',
        re.DOTALL
    )

    matches = section_regex.findall(summary_text)
    section_dict = {title.strip(): content.strip() for title, content in matches if title and content}

    if not section_dict:
        print("⚠️ No section headers matched. Using default slicing fallback.")
        section_dict = {
            "Introduction": summary_text[:500],
            "Core Contributions": summary_text[500:1000],
            "Methods": summary_text[1000:1500],
            "Findings": summary_text[1500:2000],
            "Conclusion": summary_text[2000:]
        }

    return section_dict

def generate_bullet_points_for_section(section_title, section_content, model, tokenizer, device):
    """Generate bullet points for a section."""
    prompt = f"""
You are creating bullet points for a PowerPoint slide titled \"{section_title}\" for a presentation on international business strategy.

The slide will be based on this section from an academic paper summary:
"{section_content}"

Create EXACTLY 3 bullet points that:
1. Highlight the key insights from this section
2. Are written as complete sentences (10-15 words max each)
3. Are clear, concise, and academically sound

ONLY output the 3 bullet points, one per line, no quotes or numbering.
"""
    
    response = simple_llm_call(prompt, model, tokenizer, device)
    lines = [line.strip() for line in response.split('\n') if line.strip()]

    bullet_points = []
    for line in lines:
        cleaned = re.sub(r'^[-\u2022*\d\.]+\s*', '', line)
        if cleaned:
            if len(cleaned) > 120:
                break_point = cleaned.rfind('.', 0, 120)
                if break_point > 50:
                    cleaned = cleaned[:break_point+1]
                else:
                    break_point = cleaned.rfind(' ', 90, 120)
                    if break_point > 0:
                        cleaned = cleaned[:break_point] + "."
            bullet_points.append(cleaned)

    if len(bullet_points) != 3:
        print(f"⚠️ Didn't get exactly 3 bullet points for '{section_title}'. Using fallback.")
        sentences = re.split(r'(?<=[.!?])\s+', section_content)
        bullet_points = [s.strip()[:120] for s in sentences if len(s.split()) >= 5][:3]

    while len(bullet_points) < 3:
        bullet_points.append(f"Key aspects of {section_title.lower()} in multinational strategy.")

    return bullet_points[:3]

def simple_llm_call(prompt, model, tokenizer, device, max_tokens=1000):
    inputs = tokenizer(prompt, return_tensors="pt").to(device)
    output = model.generate(
        **inputs,
        max_new_tokens=max_tokens,
        do_sample=True,
        temperature=0.7,
        top_p=0.95
    )
    return tokenizer.decode(output[0], skip_special_tokens=True).strip()

def generate_slides_from_summary(summary_text, model, tokenizer, device):
    print("\nExtracting sections from summary...")
    sections = extract_sections_from_summary(summary_text)
    print("\n✅ Extracted sections:")
    for title in sections:
        print(f" - {title}: {len(sections[title].split())} words")

    slides = []
    for title, content in sections.items():
        print(f"Generating bullet points for: {title}")
        bullet_points = generate_bullet_points_for_section(title, content, model, tokenizer, device)
        slides.append({"title": title, "bullet_points": bullet_points})

    return slides

################################################################################
# 2. PowerPoint Creation
################################################################################

def create_powerpoint(slides_data, paper_title, output_filename='presentation.pptx'):
    try:
        from pptx import Presentation
        from pptx.util import Pt

        prs = Presentation()
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = paper_title
        title.text_frame.paragraphs[0].font.size = Pt(40)
        subtitle.text = "Summary-Based Slide Deck"
        subtitle.text_frame.paragraphs[0].font.size = Pt(24)

        for slide_data in slides_data:
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text = slide_data["title"]
            title.text_frame.paragraphs[0].font.size = Pt(36)

            content = slide.placeholders[1]
            tf = content.text_frame
            tf.clear()

            for bp in slide_data["bullet_points"]:
                p = tf.add_paragraph()
                p.text = bp
                p.level = 0
                p.font.size = Pt(20)

        prs.save(output_filename)
        print(f"\n✅ PowerPoint presentation saved as '{output_filename}'")
        return True

    except Exception as e:
        print(f"\n❌ Error creating PowerPoint: {e}")
        return False

################################################################################
# 3. Main Runner
################################################################################

if __name__ == "__main__":
    folder = "/work/Chatbot-in-academia/papers-cleaned"

    print("Loading LLaMA model...")
    tokenizer, model, device = load_llama_instructor()

    print("\n📚 Summary-Based Slide Generator is ready. Type your paper title. Type 'exit' to quit.")

    while True:
        query = input("\nYou: ")
        if query.lower() == "exit":
            break

        try:
            print(f"🔍 Searching for document matching: '{query}'...")
            final_summary = llama_summarize_text_doc(query, folder, tokenizer, model, device)

            if isinstance(final_summary, str) and final_summary.startswith("I'm not sure"):
                print(f"❌ {final_summary}")
                continue

            with open("generated_summary.txt", "w", encoding="utf-8") as f:
                f.write(final_summary)

            print("✅ Final summary saved to 'generated_summary.txt'")
            paper_title = query.title()
            slides_data = generate_slides_from_summary(final_summary, model, tokenizer, device)
            create_powerpoint(slides_data, paper_title)

        except Exception as e:
            print(f"❌ An error occurred: {repr(e)}")
            print("⚠️ Try again or try a different query.")
